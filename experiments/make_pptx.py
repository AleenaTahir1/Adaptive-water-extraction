"""Generate a clean minimalist .pptx presentation.

Design principles:
  * Pure white background.
  * ONE design element per slide: a thin colored rule at the top (3 pt).
    No circles, half-circles, dot constellations, or corner motifs.
  * Generous whitespace, left-aligned typography, modern color palette.
  * Each section has its own accent color so the deck has visual rhythm
    without any individual slide being busy.
  * 16:9 widescreen, matches modern PowerPoint defaults.

Run:
    python -m experiments.make_pptx
Output:
    report/presentation.pptx
"""
from __future__ import annotations

from pathlib import Path
from typing import Sequence

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree


# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Margins / layout grid.
LEFT_MARGIN  = Inches(0.7)
RIGHT_MARGIN = Inches(0.7)
TOP_BAR_H    = Inches(0.06)            # the only decorative element

# Modern palette: charcoal primary, one accent per section.
INK         = RGBColor(0x0F, 0x17, 0x2A)  # near-black for headings
BODY_GRAY   = RGBColor(0x33, 0x3D, 0x4F)  # readable body
MUTED       = RGBColor(0x64, 0x74, 0x8B)
HAIRLINE    = RGBColor(0xE2, 0xE8, 0xF0)
TINT_GRAY   = RGBColor(0xF8, 0xFA, 0xFC)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)

# Accent colors -- one per section.
CORAL   = RGBColor(0xDC, 0x26, 0x26)
TEAL    = RGBColor(0x08, 0x91, 0xB2)
AMBER   = RGBColor(0xD9, 0x77, 0x06)
INDIGO  = RGBColor(0x4F, 0x46, 0xE5)
EMERALD = RGBColor(0x05, 0x96, 0x69)
PURPLE  = RGBColor(0x7C, 0x3A, 0xED)


PROJECT_ROOT = Path(__file__).resolve().parent.parent
FIG_DIR = PROJECT_ROOT / "figures"
OUT_PATH = PROJECT_ROOT / "report" / "presentation.pptx"


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------


def set_slide_background(slide, color: RGBColor):
    """Set solid background on a slide (overrides the master)."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h, fill_color: RGBColor,
             line_color: RGBColor | None = None, line_w_pt: float | None = None):
    """Insert a colored rectangle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
        if line_w_pt is not None:
            shape.line.width = Pt(line_w_pt)
    # Remove default shadow.
    sp = shape.shadow
    return shape


def add_textbox(slide, x, y, w, h, text: str, *,
                font_name: str = "Calibri", font_size: int = 16,
                color: RGBColor = INK, bold: bool = False,
                italic: bool = False, align: str = "left",
                line_spacing: float | None = None,
                vertical_anchor: str = "top") -> "Shape":
    """Insert a textbox with formatted single-line / paragraph text."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    if vertical_anchor == "middle":
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    elif vertical_anchor == "bottom":
        tf.vertical_anchor = MSO_ANCHOR.BOTTOM
    else:
        tf.vertical_anchor = MSO_ANCHOR.TOP

    p = tf.paragraphs[0]
    if line_spacing:
        p.line_spacing = line_spacing
    align_map = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}
    p.alignment = align_map.get(align, PP_ALIGN.LEFT)
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb


def add_section_chip(slide, x, y, label: str, color: RGBColor):
    """Tiny chip: short colored rule + uppercase label."""
    # Rule.
    rule = slide.shapes.add_connector(1, x, y + Inches(0.07), x + Inches(0.30), y + Inches(0.07))
    rule.line.color.rgb = color
    rule.line.width = Pt(2.5)
    # Label text.
    add_textbox(slide, x + Inches(0.42), y - Inches(0.03), Inches(3), Inches(0.25),
                label.upper(), font_name="Calibri", font_size=10, bold=True,
                color=color)


def add_top_bar(slide, color: RGBColor):
    """The ONLY decorative element on a content slide: a thin top color bar."""
    add_rect(slide, 0, 0, SLIDE_W, TOP_BAR_H, fill_color=color)


def add_bullets(slide, x, y, w, h, items: Sequence[str], *,
                color: RGBColor = BODY_GRAY, font_size: int = 14,
                bullet_color: RGBColor = INK):
    """Insert a clean bullet list (Calibri body, modern triangle bullet)."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = 1.30
        # Add a hanging-indent bullet via the marker character.
        bullet_run = p.add_run()
        bullet_run.text = "• "
        bullet_run.font.name = "Calibri"
        bullet_run.font.size = Pt(font_size + 2)
        bullet_run.font.color.rgb = bullet_color
        bullet_run.font.bold = True
        # Body text.
        body_run = p.add_run()
        body_run.text = item
        body_run.font.name = "Calibri"
        body_run.font.size = Pt(font_size)
        body_run.font.color.rgb = color
        # Space between bullets.
        p.space_after = Pt(8)
    return tb


def add_table(slide, x, y, w, h, data: list[list[str]], col_fracs: list[float],
              header_color: RGBColor, alt_color: RGBColor = TINT_GRAY,
              header_font_color: RGBColor = WHITE,
              header_font_size: int = 12, body_font_size: int = 11):
    """Insert a styled table. col_fracs sums to 1.0."""
    rows, cols = len(data), len(data[0])
    table_shape = slide.shapes.add_table(rows, cols, x, y, w, h)
    table = table_shape.table
    # Column widths.
    for j, frac in enumerate(col_fracs):
        table.columns[j].width = Emu(int(w * frac))
    # Style cells.
    for i, row in enumerate(data):
        for j, cell_text in enumerate(row):
            cell = table.cell(i, j)
            cell.margin_left = Inches(0.10)
            cell.margin_right = Inches(0.10)
            cell.margin_top = Inches(0.05)
            cell.margin_bottom = Inches(0.05)
            # Cell fill.
            cell.fill.solid()
            if i == 0:
                cell.fill.fore_color.rgb = header_color
            else:
                cell.fill.fore_color.rgb = WHITE if i % 2 == 1 else alt_color
            # Cell text.
            tf = cell.text_frame
            tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = cell_text
            run.font.name = "Calibri"
            if i == 0:
                run.font.size = Pt(header_font_size)
                run.font.bold = True
                run.font.color.rgb = header_font_color
            else:
                run.font.size = Pt(body_font_size)
                run.font.color.rgb = INK
                run.font.bold = (j == 0)  # bold first column
    return table_shape


def add_image_centered(slide, path: Path, x, y, max_w, max_h):
    """Add image, centered horizontally in the [x, x+max_w] box, anchored at y."""
    if not path.exists():
        return None
    pic = slide.shapes.add_picture(str(path), x, y, width=max_w, height=max_h)
    # python-pptx scales independently; correct to maintain aspect ratio.
    iw_native = pic.image.size[0]
    ih_native = pic.image.size[1]
    aspect = iw_native / ih_native
    target_w = max_w
    target_h = int(target_w / aspect)
    if target_h > max_h:
        target_h = max_h
        target_w = int(target_h * aspect)
    # Re-set sizes to preserve aspect.
    pic.width = Emu(target_w)
    pic.height = Emu(target_h)
    # Center horizontally.
    pic.left = Emu(int(x + (max_w - target_w) / 2))
    pic.top = Emu(y)
    return pic


def add_slide_head(slide, chip_label: str, title: str, subtitle: str | None,
                   color: RGBColor):
    """Standard slide head: top bar + chip + title (+ optional subtitle)."""
    add_top_bar(slide, color)
    chip_y = Inches(0.45)
    add_section_chip(slide, LEFT_MARGIN, chip_y, chip_label, color)
    add_textbox(slide, LEFT_MARGIN, Inches(0.85), SLIDE_W - LEFT_MARGIN - RIGHT_MARGIN, Inches(0.9),
                title, font_name="Calibri", font_size=34, bold=True, color=INK)
    if subtitle:
        add_textbox(slide, LEFT_MARGIN, Inches(1.55),
                    SLIDE_W - LEFT_MARGIN - RIGHT_MARGIN, Inches(0.4),
                    subtitle, font_name="Calibri Light", font_size=14,
                    italic=True, color=MUTED)


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------


def new_blank_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 6 = blank
    set_slide_background(slide, WHITE)
    return slide


def slide_01_title(prs):
    slide = new_blank_slide(prs)
    # Top bar in indigo.
    add_top_bar(slide, INDIGO)

    # Left-aligned title block.
    tx = Inches(0.9)

    # Small "01 / TOPIC" mark.
    add_textbox(slide, tx, Inches(2.0), Inches(2.5), Inches(0.3),
                "01  TOPIC", font_size=11, bold=True, color=INDIGO,
                font_name="Calibri")

    # Topic name -- two-line big serif.
    add_textbox(slide, tx, Inches(2.55), Inches(11), Inches(1.5),
                "Adaptive Water Extraction", font_name="Cambria",
                font_size=56, bold=True, color=INK)
    add_textbox(slide, tx, Inches(3.55), Inches(11), Inches(1.5),
                "with Q Learning under Climate Stress",
                font_name="Cambria", font_size=30, italic=True,
                color=BODY_GRAY)

    # Thin horizontal accent rule.
    rule = slide.shapes.add_connector(1, tx, Inches(4.7),
                                       tx + Inches(3.5), Inches(4.7))
    rule.line.color.rgb = INDIGO
    rule.line.width = Pt(2)

    add_textbox(slide, tx, Inches(4.85), Inches(8), Inches(0.4),
                "An agent based model in Python.  Mesa framework.",
                font_size=14, italic=True, color=MUTED,
                font_name="Calibri Light")


def slide_02_team(prs):
    slide = new_blank_slide(prs)
    add_slide_head(slide, "team", "Group Members", None, TEAL)

    # Five rows, simple and clean.
    members = [
        ("01", "Saqlain Abbas"),
        ("02", "Aleena Tahir"),
        ("03", "Aena Habib"),
        ("04", "Eman Asghar"),
        ("05", "Dua Kamal"),
    ]
    row_h = Inches(0.7)
    start_y = Inches(2.6)
    for i, (num, name) in enumerate(members):
        y = start_y + row_h * i
        # Serial number in teal.
        add_textbox(slide, Inches(2.0), y, Inches(0.7), Inches(0.6),
                    num, font_name="Cambria", font_size=22, bold=True,
                    color=TEAL, vertical_anchor="middle")
        # Vertical separator.
        sep = slide.shapes.add_connector(
            1, Inches(2.9), y + Inches(0.12), Inches(2.9), y + Inches(0.48))
        sep.line.color.rgb = HAIRLINE
        sep.line.width = Pt(0.75)
        # Name.
        add_textbox(slide, Inches(3.1), y, Inches(8), Inches(0.6),
                    name, font_name="Calibri", font_size=24, color=INK,
                    vertical_anchor="middle")


def slide_03_contents(prs):
    slide = new_blank_slide(prs)
    add_slide_head(slide, "agenda", "Contents", None, AMBER)

    items = [
        ("01", "The Problem and the Research Question"),
        ("02", "How the Model Works"),
        ("03", "Climate Scenarios and Experimental Design"),
        ("04", "Result 1.  Tragedy Baseline"),
        ("05", "Result 2.  Enforcement Saves the Stable Climate"),
        ("06", "Result 3.  Why Enforcement Fails Under Stress"),
        ("07", "Primary RQ.  Effect of Discount Factor"),
        ("08", "Sensitivity and Validation"),
        ("09", "Conclusions and Future Work"),
    ]
    row_h = Inches(0.42)
    start_y = Inches(2.3)
    for i, (num, title) in enumerate(items):
        y = start_y + row_h * i
        add_textbox(slide, Inches(0.9), y, Inches(0.8), Inches(0.4),
                    num, font_name="Cambria", font_size=18, bold=True,
                    color=AMBER, vertical_anchor="middle")
        add_textbox(slide, Inches(1.7), y, Inches(11), Inches(0.4),
                    title, font_name="Calibri", font_size=14, color=INK,
                    vertical_anchor="middle")


def slide_04_problem(prs):
    slide = new_blank_slide(prs)
    add_slide_head(slide, "01 . context", "The Problem",
                   "Tragedy of the Commons meets a changing climate", CORAL)
    bullets = [
        "Twenty farmers share one river. Each season every farmer decides how much water to extract.",
        "If everyone extracts heavily, the resource collapses and everyone loses. Hardin (1968) called this the Tragedy of the Commons.",
        "Pakistan's Indus Basin supports over 90 percent of national agriculture and is one of the world's most water stressed basins.",
        "By 2030 nearly half the global population will live in regions of high water stress (UN World Water Development Report).",
        "Real farmers learn, observe neighbors, and adapt. Existing agent based models do not capture this learning.",
    ]
    add_bullets(slide, LEFT_MARGIN, Inches(2.2),
                SLIDE_W - LEFT_MARGIN - RIGHT_MARGIN, Inches(4.5),
                bullets, font_size=15, bullet_color=CORAL)


def slide_05_research_question(prs):
    slide = new_blank_slide(prs)
    add_slide_head(slide, "01 . focus", "Research Question", None, PURPLE)

    # Primary question card with subtle left rule.
    card_x = LEFT_MARGIN
    card_y = Inches(2.0)
    card_w = SLIDE_W - LEFT_MARGIN - RIGHT_MARGIN
    card_h = Inches(1.7)
    # Tinted background.
    bg = add_rect(slide, card_x, card_y, card_w, card_h,
                  RGBColor(0xF5, 0xF3, 0xFF))
    # Left accent bar.
    add_rect(slide, card_x, card_y, Inches(0.10), card_h, PURPLE)

    add_textbox(slide, card_x + Inches(0.3), card_y + Inches(0.2),
                Inches(3), Inches(0.3),
                "PRIMARY QUESTION", font_size=10, bold=True, color=PURPLE)
    add_textbox(slide, card_x + Inches(0.3), card_y + Inches(0.55),
                card_w - Inches(0.6), Inches(1.1),
                "How does the discount factor gamma in Q learning agents "
                "affect the emergence of sustainable water extraction "
                "strategies under varying climate stress scenarios?",
                font_name="Cambria", font_size=18, italic=True, color=INK,
                line_spacing=1.25)

    # Sub-questions.
    add_textbox(slide, LEFT_MARGIN, Inches(4.0), Inches(3), Inches(0.3),
                "SUB QUESTIONS", font_size=10, bold=True, color=PURPLE)
    subs = [
        "Can Q learning agents discover cooperative strategies without external enforcement?",
        "How do learned strategies break down under sudden drought shocks, and how fast do agents re adapt?",
        "At what monitoring intensity does enforcement become unnecessary?",
    ]
    add_bullets(slide, LEFT_MARGIN, Inches(4.35),
                SLIDE_W - LEFT_MARGIN - RIGHT_MARGIN, Inches(2.2),
                subs, font_size=14, bullet_color=PURPLE)


def slide_06_model(prs):
    slide = new_blank_slide(prs)
    add_slide_head(slide, "02 . model", "How the Model Works",
                   "Three stage step.  Act, monitor, regenerate, finalize.", INDIGO)
    bullets = [
        "Each irrigator chooses one of five extraction levels: 0, 25, 50, 75, 100 percent of max.",
        "Observed state. Water level bin, own previous action, neighbors average action, climate stress flag.",
        "Reward equals log utility of water extracted minus a sustainability penalty minus any fine.",
        "Shared resource follows logistic regeneration with a baseline inflow modulated by climate factor.",
        "Monitor agent catches over extractors with probability p detect and issues a fine.",
        "Social network is a random k regular graph. Each farmer observes k neighbors.",
        "Implemented in Python 3.12 with Mesa 2.3.4. About 700 lines of clear code.",
    ]
    add_bullets(slide, LEFT_MARGIN, Inches(2.2),
                SLIDE_W - LEFT_MARGIN - RIGHT_MARGIN, Inches(4.5),
                bullets, font_size=14, bullet_color=INDIGO)


def slide_07_q_learning(prs):
    slide = new_blank_slide(prs)
    add_slide_head(slide, "02 . math", "The Q Learning Update Rule",
                   "Watkins and Dayan, 1992.  Tabular and interpretable.", PURPLE)

    # Formula card.
    card_y = Inches(2.4)
    card_h = Inches(1.6)
    card_w = SLIDE_W - LEFT_MARGIN - RIGHT_MARGIN
    add_rect(slide, LEFT_MARGIN, card_y, card_w, card_h,
             RGBColor(0xF5, 0xF3, 0xFF))
    add_rect(slide, LEFT_MARGIN, card_y, Inches(0.10), card_h, PURPLE)

    # Formula text.
    add_textbox(slide, LEFT_MARGIN + Inches(0.3), card_y + Inches(0.3),
                card_w - Inches(0.6), Inches(0.6),
                "Q(s, a)   =   Q(s, a)   +   alpha [ r + gamma . max Q(s', a')   minus   Q(s, a) ]",
                font_name="Cambria", font_size=22, bold=True, color=INK,
                align="center")
    add_textbox(slide, LEFT_MARGIN + Inches(0.3), card_y + Inches(1.05),
                card_w - Inches(0.6), Inches(0.4),
                "alpha is the learning rate.   gamma is the discount factor.   s' is the post regen state.",
                font_name="Cambria", font_size=12, italic=True, color=MUTED,
                align="center")

    bullets = [
        "Exploration uses epsilon greedy with exponential decay. Epsilon starts at 0.30, decays at rate 0.010.",
        "State space is small. About 150 unique state action pairs. Fits in a tabular Q table.",
        "Action 2 equals each agent's fair share of total sustainable yield. Auto scaled with N.",
        "Heterogeneity. Each agent has its own RNG and its own initial Q table noise.",
    ]
    add_bullets(slide, LEFT_MARGIN, Inches(4.4),
                SLIDE_W - LEFT_MARGIN - RIGHT_MARGIN, Inches(2.5),
                bullets, font_size=12, bullet_color=PURPLE)


def slide_08_climate(prs):
    slide = new_blank_slide(prs)
    add_slide_head(slide, "03 . scenarios", "Climate Scenarios",
                   "Four climate factor profiles test resilience.", TEAL)
    data = [
        ["Scenario",   "Climate factor profile"],
        ["Stable",     "c_f = 1.0 forever.  Baseline.  No climate change."],
        ["Gradual",    "c_f declines linearly from 1.0 to 0.5 over 400 steps."],
        ["Shock",      "c_f = 1.0 normally.  Drops to 0.3 between steps 200 and 300."],
        ["Stochastic", "c_f is gaussian noise around 0.85, clipped to [0.3, 1.0]."],
    ]
    add_table(slide, LEFT_MARGIN, Inches(2.3),
              SLIDE_W - LEFT_MARGIN - RIGHT_MARGIN, Inches(3),
              data, [0.22, 0.78], header_color=TEAL)

    add_textbox(slide, LEFT_MARGIN, Inches(6.0),
                SLIDE_W - LEFT_MARGIN - RIGHT_MARGIN, Inches(0.3),
                "c_f multiplies BOTH logistic regeneration AND baseline inflow.",
                font_size=12, italic=True, color=MUTED)
    add_textbox(slide, LEFT_MARGIN, Inches(6.3),
                SLIDE_W - LEFT_MARGIN - RIGHT_MARGIN, Inches(0.3),
                "Lower c_f means less water replenished each season.",
                font_size=12, italic=True, color=MUTED)


def slide_09_experiments(prs):
    slide = new_blank_slide(prs)
    add_slide_head(slide, "03 . experiments", "Experimental Design",
                   "Three batch experiments totaling 340 runs.", AMBER)
    data = [
        ["Experiment",              "Variables",                                   "Runs"],
        ["A.  Tragedy baseline",    "4 climates x 30 seeds.  No monitor.",         "120"],
        ["B.  With enforcement",    "4 climates x 30 seeds.  p detect = 0.3",      "120"],
        ["C.  Gamma sweep",         "stable and shock x 5 gamma x 10 seeds",       "100"],
    ]
    add_table(slide, LEFT_MARGIN, Inches(2.3),
              SLIDE_W - LEFT_MARGIN - RIGHT_MARGIN, Inches(2.4),
              data, [0.30, 0.55, 0.15], header_color=AMBER)

    extras = [
        "Each run is 500 simulation steps.  About 125 seasons if one step is one quarter year.",
        "Sensitivity analysis varies alpha, gamma, epsilon, regeneration rate, p detect, N by plus or minus 20 percent.",
        "Validation.  Three sanity checks against Hardin 1968, Perolat 2017, Indus Basin 2022.",
    ]
    add_bullets(slide, LEFT_MARGIN, Inches(5.0),
                SLIDE_W - LEFT_MARGIN - RIGHT_MARGIN, Inches(2.0),
                extras, font_size=13, bullet_color=AMBER)


def slide_with_fig(prs, chip: str, title: str, subtitle: str, fig: str,
                   takeaways: list[str], color: RGBColor):
    slide = new_blank_slide(prs)
    add_slide_head(slide, chip, title, subtitle, color)

    # Figure on left.
    fig_x = LEFT_MARGIN
    fig_y = Inches(2.2)
    fig_max_w = Inches(7.2)
    fig_max_h = Inches(4.8)
    add_image_centered(slide, FIG_DIR / fig, fig_x, fig_y, fig_max_w, fig_max_h)

    # Takeaways on right.
    rx = Inches(8.3)
    rw = SLIDE_W - rx - RIGHT_MARGIN
    add_textbox(slide, rx, Inches(2.2), rw, Inches(0.3),
                "WHAT IT MEANS", font_size=10, bold=True, color=color)
    add_bullets(slide, rx, Inches(2.6), rw, Inches(4.5),
                takeaways, font_size=12, bullet_color=color)


def slide_enforcement_fails(prs):
    slide = new_blank_slide(prs)
    add_slide_head(slide, "06 . key finding", "Why Enforcement Fails Under Stress",
                   "The fixed quota policy breaks when climate changes.", CORAL)

    # Coral callout strip.
    callout_y = Inches(2.4)
    callout_h = Inches(0.55)
    add_rect(slide, LEFT_MARGIN, callout_y,
             SLIDE_W - LEFT_MARGIN - RIGHT_MARGIN, callout_h, CORAL)
    add_textbox(slide, LEFT_MARGIN + Inches(0.25), callout_y,
                SLIDE_W - LEFT_MARGIN - RIGHT_MARGIN - Inches(0.5), callout_h,
                "Even p detect = 0.9 cannot save the drought shock scenario.",
                font_name="Cambria", font_size=18, bold=True, color=WHITE,
                vertical_anchor="middle")

    # Evidence table.
    data = [
        ["p detect", "stable.  final water", "shock.  final water"],
        ["0.3", "700", "5.6  (collapsed)"],
        ["0.6", "777", "5.6  (collapsed)"],
        ["0.9", "787", "5.6  (collapsed)"],
    ]
    add_table(slide, LEFT_MARGIN, Inches(3.1),
              SLIDE_W - LEFT_MARGIN - RIGHT_MARGIN, Inches(1.8),
              data, [0.20, 0.40, 0.40], header_color=CORAL,
              alt_color=RGBColor(0xFE, 0xF2, 0xF2))

    # Why it happens.
    add_textbox(slide, LEFT_MARGIN, Inches(5.2),
                SLIDE_W - LEFT_MARGIN - RIGHT_MARGIN, Inches(0.3),
                "WHY IT HAPPENS  (structural, not a tuning issue)",
                font_size=10, bold=True, color=CORAL)
    reasons = [
        "Agents' five action grid is calibrated to baseline climate where c f = 1.0",
        "When c f drops to 0.3 during the shock, even the cooperative fair share action is over extraction.",
        "Sixty eight percent of agents are STILL choosing cooperative actions during failure.  Not selfish.",
        "Implication.  Real regulators need adaptive (climate aware) quotas.  Fixed quotas are not enough.",
    ]
    add_bullets(slide, LEFT_MARGIN, Inches(5.55),
                SLIDE_W - LEFT_MARGIN - RIGHT_MARGIN, Inches(1.8),
                reasons, font_size=12, bullet_color=CORAL)


def slide_validation(prs):
    slide = new_blank_slide(prs)
    add_slide_head(slide, "08 . validation", "Validation Against Literature",
                   "Three sanity checks.  All three pass.", EMERALD)
    data = [
        ["#", "Source",            "Test",                                                                                                  "Result"],
        ["1", "Hardin 1968",       "No enforcement should collapse.  With enforcement under stable climate it should sustain.",            "PASS  5.6 vs 504"],
        ["2", "Perolat et al 2017","Scarcity should drive aggression.  Compare frac defect steps 200-300, stable vs shock with enforcement.","PASS  25.3% vs 36.5%"],
        ["3", "Indus Basin 2022",  "A 70% drop in c f should leave a measurable dip in water level at step 300.",                          "PASS  507 vs 1.5"],
    ]
    add_table(slide, LEFT_MARGIN, Inches(2.2),
              SLIDE_W - LEFT_MARGIN - RIGHT_MARGIN, Inches(2.2),
              data, [0.05, 0.17, 0.53, 0.25], header_color=EMERALD,
              alt_color=RGBColor(0xEC, 0xFD, 0xF5))

    add_image_centered(slide, FIG_DIR / "fig07_validation.png",
                       Inches(2.0), Inches(4.7), Inches(9.3), Inches(2.5))


def slide_spatial(prs):
    slide = new_blank_slide(prs)
    add_slide_head(slide, "08 . evidence", "Emergent Spatial Dynamics",
                   "Three snapshots from one shock scenario run.", AMBER)

    img_w = Inches(3.9)
    img_h = Inches(4.0)
    img_y = Inches(2.2)
    captions = ["Step 199.  Pre shock cooperation",
                "Step 260.  Mid shock pressure",
                "Step 450.  Post shock equilibrium"]
    files = ["fig08_grid_pre_shock.png",
             "fig09_grid_during_shock.png",
             "fig10_grid_recovery.png"]
    gap = Inches(0.3)
    total_w = img_w * 3 + gap * 2
    x_start = (SLIDE_W - total_w) // 2
    for i, (cap, f) in enumerate(zip(captions, files)):
        x = x_start + (img_w + gap) * i
        add_image_centered(slide, FIG_DIR / f, x, img_y, img_w, img_h)
        add_textbox(slide, x, img_y + img_h + Inches(0.1),
                    img_w, Inches(0.4),
                    cap, font_size=12, bold=True, color=AMBER, align="center",
                    font_name="Cambria")


def slide_conclusion(prs):
    slide = new_blank_slide(prs)
    add_slide_head(slide, "09 . wrap up", "Conclusions and Future Work",
                   "What we learned. What comes next.", INDIGO)

    # Findings header
    add_textbox(slide, LEFT_MARGIN, Inches(2.2), Inches(3), Inches(0.3),
                "FINDINGS", font_size=10, bold=True, color=INDIGO)
    findings = [
        "Pure Q learners reproduce the tragedy of the commons in every climate scenario.",
        "Moderate enforcement at p detect = 0.3 saves cooperation under stable climate.  Water rises from 5 to 685.",
        "The SAME enforcement fails under any climate stress.  Even p detect = 0.9 cannot rescue shock.",
        "Gamma matters at the margins under stable climate but is dominated by climate dynamics under stress.",
    ]
    add_bullets(slide, LEFT_MARGIN, Inches(2.55),
                SLIDE_W - LEFT_MARGIN - RIGHT_MARGIN, Inches(2.4),
                findings, font_size=13, bullet_color=INDIGO)

    add_textbox(slide, LEFT_MARGIN, Inches(5.0), Inches(3), Inches(0.3),
                "FUTURE WORK", font_size=10, bold=True, color=CORAL)
    future = [
        "Continuous or adaptive action space so agents can scale extraction with current regeneration.",
        "Climate aware monitor that tightens the fair share threshold during droughts.",
        "Deep reinforcement learning extension.  Comparison with human commons experiments.",
    ]
    add_bullets(slide, LEFT_MARGIN, Inches(5.35),
                SLIDE_W - LEFT_MARGIN - RIGHT_MARGIN, Inches(2.0),
                future, font_size=13, bullet_color=CORAL)


def slide_thank_you(prs):
    slide = new_blank_slide(prs)
    # Single thin indigo top bar -- and nothing else decorative.
    add_top_bar(slide, INDIGO)

    add_textbox(slide, Inches(0), Inches(2.8), SLIDE_W, Inches(2.0),
                "Thank you", font_name="Cambria", font_size=100, bold=True,
                color=INK, align="center")
    # Subtle thin rule under the headline.
    rule = slide.shapes.add_connector(
        1, Inches(4.5), Inches(5.0), Inches(8.8), Inches(5.0))
    rule.line.color.rgb = INDIGO
    rule.line.width = Pt(2)


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------


def build_pptx():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_01_title(prs)
    slide_02_team(prs)
    slide_03_contents(prs)
    slide_04_problem(prs)
    slide_05_research_question(prs)
    slide_06_model(prs)
    slide_07_q_learning(prs)
    slide_08_climate(prs)
    slide_09_experiments(prs)

    slide_with_fig(prs,
        chip="04 . result 1",
        title="Tragedy Baseline",
        subtitle="Without enforcement every climate collapses.",
        fig="fig01_water_by_scenario.png",
        takeaways=[
            "All four scenarios collapse to roughly 5 units of water within 80 steps.",
            "Mean payoff falls to between minus 755 and minus 804 per agent.",
            "Standard deviation across 30 seeds is effectively zero.  Tragedy is deterministic.",
            "Reproduces Hardin (1968) at the population level.",
        ],
        color=AMBER)

    slide_with_fig(prs,
        chip="05 . result 2",
        title="Enforcement Saves the Stable Climate",
        subtitle="Moderate monitoring flips the dynamics.",
        fig="fig02_enforcement_effect.png",
        takeaways=[
            "Stable scenario.  Water rises from 5 to 685.  A 125 times improvement.",
            "Mean payoff turns positive.  Minus 756 becomes plus 222 per agent.",
            "Eighty six percent of agents choose cooperative actions.",
            "Three other scenarios still collapse.  See next slide for why.",
        ],
        color=EMERALD)

    slide_enforcement_fails(prs)

    slide_with_fig(prs,
        chip="07 . primary RQ",
        title="Effect of the Discount Factor",
        subtitle="Gamma sweep over stable and shock scenarios.",
        fig="fig03_gamma_sweep.png",
        takeaways=[
            "Stable scenario.  Final water mildly DECREASES with gamma.  From 748 down to 688.",
            "Shock scenario.  Final water is identical at 5.55 across all gamma values.",
            "Surprising direction.  Lower gamma is slightly better in stable climate.",
            "Refines the answer.  Gamma is not the resilience knob.  Action space matters more.",
        ],
        color=PURPLE)

    slide_with_fig(prs,
        chip="08 . sensitivity",
        title="Sensitivity Analysis",
        subtitle="Plus or minus 20 percent on six parameters.",
        fig="fig04_sensitivity_tornado.png",
        takeaways=[
            "p detect dominates.  Plus or minus 38 units of final water.",
            "Epsilon zero and alpha are non monotonic.  Baseline near local optimum.",
            "Gamma effect moderate.  Regeneration rate small.",
            "Validates that p detect = 0.3 is a fair choice for headline runs.",
        ],
        color=TEAL)

    slide_validation(prs)
    slide_spatial(prs)
    slide_conclusion(prs)
    slide_thank_you(prs)

    prs.save(str(OUT_PATH))
    print(f"Wrote {OUT_PATH}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    build_pptx()
